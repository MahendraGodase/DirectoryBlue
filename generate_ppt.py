from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(31, 78, 121)
TEXT_COLOR = RGBColor(64, 64, 64)

def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for item in content:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "DirectoryBlue"
slide.placeholders[1].text = "Angular Application with Azure SQL & MCP Integration\n\nPresented by: Bob - AI Software Engineer"

# Slide 2
add_slide("Project Overview", [
    "Purpose: Provider information management",
    "Technology: Angular 20.3.0 with TypeScript",
    "Backend: Azure SQL Database + Mock Server",
    "Integration: IBM Context Studio via MCP",
    "Design: Responsive, modern UI"
])

# Slide 3
add_slide("Key Features", [
    "Search Inventory Management",
    "Import File Management",
    "Multiple Inventory Views",
    "Bulk Operations",
    "Export to Excel"
])

# Slide 4
add_slide("Technology Stack", [
    "Frontend: Angular 20.3.0, TypeScript 5.9.2",
    "Backend: Node.js 24.11.1, Express.js",
    "Database: Azure SQL Database",
    "Integration: MCP, JIRA",
    "Tools: VS Code, Cline, Git"
])

# Slide 5
add_slide("MCP Integration", [
    "Model Context Protocol (MCP)",
    "Connects Cline to IBM Context Studio",
    "Standardized protocol for AI tools",
    "Secure authentication with JWT",
    "API Key: ctx_b7893df16d01"
])

# Slide 6
add_slide("MCP Configuration", [
    "Server Name: context-studio",
    "Type: streamable-http",
    "URL: IBM Context Studio Gateway",
    "Authentication: Bearer Token + API Key",
    "Status: Enabled"
])

# Slide 7
add_slide("MCP Permissions", [
    "gateways.read - Read gateway info",
    "servers.use - Use MCP servers",
    "tools.execute - Execute tools",
    "resources.read - Access resources",
    "Token Valid: Jan 2025 - Jun 2026"
])

# Slide 8
add_slide("Prerequisites", [
    "Node.js v18.19 or higher",
    "npm v9.0 or higher",
    "Angular CLI v20.3.10",
    "VS Code with Cline extension",
    "PowerShell (Windows)"
])

# Slide 9
add_slide("Installation Steps", [
    "1. Install backend dependencies:",
    "   cd mock-backend && npm install",
    "",
    "2. Install frontend dependencies:",
    "   cd DirectoryBlue && npm install"
])

# Slide 10
add_slide("Starting the Application", [
    "Backend Server (Terminal 1):",
    "  cd mock-backend && npm start",
    "  Runs on: http://localhost:3000",
    "",
    "Frontend App (Terminal 2):",
    "  npm start",
    "  Runs on: http://localhost:4200"
])

# Slide 11
add_slide("Application Interface", [
    "Navigation: Home, Admin, Reports, Timekeeper",
    "Welcome message with user name",
    "Tabs: Search, Import, My Inventory",
    "Survey & CAQH Unassigned Inventory",
    "Search parameters and filters"
])

# Slide 12
add_slide("Search Functionality", [
    "Search by NPI (National Provider ID)",
    "Filter by Task ID",
    "Data Source: CAQH, Survey, Manual",
    "Date range filtering",
    "Task status filtering",
    "Clear filter option"
])

# Slide 13
add_slide("Mock Test Data", [
    "TSK001 - Dr. John Smith (In Progress)",
    "TSK002 - Dr. Jane Doe (Pending)",
    "TSK003 - Dr. Robert Johnson (Completed)",
    "TSK004 - Dr. Sarah Williams (In Progress)",
    "TSK005 - Dr. Michael Brown (Cancelled)"
])

# Slide 14
add_slide("Testing Scenarios", [
    "1. Search all tasks (empty search)",
    "2. Search by NPI: 1234567890",
    "3. Filter by Data Source: CAQH",
    "4. Filter by Status: In Progress",
    "5. Clear all filters"
])

# Slide 15
add_slide("Bulk Operations", [
    "Select individual tasks with checkboxes",
    "Select all with header checkbox",
    "Bulk Update multiple tasks",
    "Export selected data to Excel",
    "View success messages"
])

# Slide 16
add_slide("Project Structure", [
    "src/app/ - Application components",
    "src/app/services/ - API services",
    "src/app/models/ - Data models",
    "mock-backend/ - Backend server",
    "src/environments/ - Configuration",
    ".bob/mcp.json - MCP settings"
])

# Slide 17
add_slide("Routing Configuration", [
    "/ → Redirect to /home",
    "/home → HomeComponent",
    "/about → AboutComponent",
    "/contact → ContactComponent",
    "** → 404 handler (redirect to home)"
])

# Slide 18
add_slide("Environment Configuration", [
    "Development (environment.ts):",
    "  apiEndpoint: http://localhost:3000/api",
    "",
    "Production (environment.prod.ts):",
    "  apiEndpoint: Azure Functions URL",
    "  production: true"
])

# Slide 19
add_slide("Azure SQL Integration", [
    "Server: Azure SQL Server",
    "Database: ProviderInfo",
    "Tables: OutreachDetails, Tasks, Providers",
    "Authentication: Azure AD / SQL Auth",
    "Connection pooling enabled"
])

# Slide 20
add_slide("JIRA Integration", [
    "Create JIRA issues from tasks",
    "Sync task status with JIRA",
    "Track progress in JIRA",
    "Link tasks to JIRA tickets",
    "Bidirectional status sync"
])

# Slide 21
add_slide("API Endpoints", [
    "GET /api/tasks - Get all tasks",
    "GET /api/tasks/:id - Get specific task",
    "POST /api/tasks/search - Search with filters",
    "POST /api/tasks - Create new task",
    "PUT /api/tasks/:id - Update task",
    "POST /api/tasks/bulk-update - Bulk operations"
])

# Slide 22
add_slide("Security Features", [
    "Azure AD integration ready",
    "JWT token support",
    "Role-based access control",
    "HTTPS enforcement",
    "API key validation",
    "SQL injection prevention"
])

# Slide 23
add_slide("Troubleshooting - Common Issues", [
    "Port in use: Kill process or change port",
    "Module not found: npm install",
    "CORS errors: Restart both servers",
    "PowerShell policy: Set-ExecutionPolicy",
    "Build errors: Clear cache and rebuild"
])

# Slide 24
add_slide("Troubleshooting - MCP", [
    "Server not showing: Verify JSON syntax",
    "401 Unauthorized: Token expired",
    "403 Forbidden: Invalid API key",
    "Network Error: Check VPN/network",
    "Solution: Restart VS Code"
])

# Slide 25
add_slide("Development Workflow", [
    "1. Start backend: cd mock-backend && npm start",
    "2. Start frontend: npm start",
    "3. Make changes (auto-reload enabled)",
    "4. Test in browser",
    "5. Check console for errors",
    "6. Commit code to Git"
])

# Slide 26
add_slide("Building for Production", [
    "Command: npm run build",
    "Output: dist/directory-blue/",
    "Optimized bundles",
    "Minified code",
    "Deploy to: Azure, AWS, Netlify"
])

# Slide 27
add_slide("Performance Optimization", [
    "Frontend: Lazy loading, code splitting",
    "Backend: Connection pooling, caching",
    "Database: Query optimization",
    "Monitoring: Application Insights",
    "Compression: Response compression"
])

# Slide 28
add_slide("Future Enhancements", [
    "Phase 1: Advanced filtering",
    "Phase 2: Full Azure SQL integration",
    "Phase 3: Mobile app",
    "AI-powered insights",
    "Real-time notifications",
    "Advanced analytics"
])

# Slide 29
add_slide("Testing Strategy", [
    "Unit Tests: npm test (Jasmine, Karma)",
    "E2E Tests: ng e2e (Protractor)",
    "Manual Testing: Browser compatibility",
    "Performance: Lighthouse",
    "Accessibility: WCAG compliance"
])

# Slide 30
add_slide("Documentation Resources", [
    "README.md - Quick start guide",
    "PROJECT_DOCUMENTATION.md - Detailed docs",
    "START_APPLICATION.md - Setup guide",
    "MCP_CONFIGURATION_GUIDE.md - MCP setup",
    "Angular Docs: https://angular.dev"
])

# Slide 31
add_slide("Best Practices", [
    "TypeScript strict mode",
    "ESLint configuration",
    "Prettier formatting",
    "Code reviews",
    "Git workflow with branches",
    "Semantic commits"
])

# Slide 32
add_slide("Team Collaboration", [
    "Version Control: Git + GitHub",
    "Communication: Daily standups",
    "Tools: VS Code, Cline, JIRA",
    "Knowledge Sharing: Documentation",
    "Code Reviews: Pull requests"
])

# Slide 33
add_slide("Monitoring & Maintenance", [
    "Application Insights",
    "Error tracking (Sentry)",
    "Performance metrics",
    "Daily database backups",
    "Security patches",
    "Dependency updates"
])

# Slide 34
add_slide("Quick Reference", [
    "Frontend: http://localhost:4200",
    "Backend: http://localhost:3000",
    "Health: http://localhost:3000/api/health",
    "Start: npm start",
    "Build: npm run build",
    "Test: npm test"
])

# Slide 35
add_slide("Summary", [
    "✅ DirectoryBlue application overview",
    "✅ MCP integration with IBM Context Studio",
    "✅ Complete setup and installation",
    "✅ Running and testing the application",
    "✅ Architecture and features",
    "✅ Best practices and troubleshooting"
])

# Slide 36
add_slide("Next Steps", [
    "1. Application is running successfully",
    "2. Test all features thoroughly",
    "3. Configure Azure SQL for production",
    "4. Set up CI/CD pipeline",
    "5. Deploy to production environment"
])

# Final Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Thank You!"
slide.placeholders[1].text = "DirectoryBlue - Ready for Production\n\n🚀 Frontend: http://localhost:4200\n📡 Backend: http://localhost:3000\n🔌 MCP: Connected to IBM Context Studio\n\nHappy Coding! 💻"

prs.save('DirectoryBlue_Presentation.pptx')
print("PowerPoint presentation created: DirectoryBlue_Presentation.pptx")

# Made with Bob
